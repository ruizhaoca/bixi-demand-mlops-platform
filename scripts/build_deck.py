"""Assemble the final presentation deck: docs/presentation/bixi_mlops_deck.pptx.

16:9 PowerPoint built with python-pptx. Strictly organised by presenter
(no interleaving): Title (shared) -> Sarah (intro + data) -> Louis (AWS &
structure) -> Othmane (modelling) -> Rui (rebalancing + demo + close) ->
Thank-you (shared). Architecture, pipeline, repo-tree and model-selection
diagrams are native, editable PPTX shapes. Per-slide speaker notes are injected
into the notes pane.

Run:  ./.venv/bin/python scripts/build_deck.py
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
PRES = REPO / "docs" / "presentation"
CHARTS = PRES / "charts"
OUT = PRES / "bixi_mlops_deck.pptx"

# --------------------------------------------------------------------------- #
# Palette & presenters
# --------------------------------------------------------------------------- #
RED = RGBColor(0xB9, 0x1C, 0x1C)
RED_DK = RGBColor(0x7F, 0x1D, 0x1D)
INK = RGBColor(0x11, 0x18, 0x27)
DARK = RGBColor(0x1F, 0x29, 0x37)
GREY = RGBColor(0x6B, 0x72, 0x80)
LGREY = RGBColor(0x9C, 0xA3, 0xAF)
PALE = RGBColor(0xF3, 0xF4, 0xF6)
PALER = RGBColor(0xF9, 0xFA, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)
TEAL = RGBColor(0x0F, 0x76, 0x6E)
AMBER = RGBColor(0xB4, 0x53, 0x09)
PURPLE = RGBColor(0x6D, 0x28, 0xD9)

FONT = "Calibri"
MONO = "Consolas"

PRESENTER = {
    "shared": ("The Bixi Crew", RED),
    "sarah": ("Sarah Liu", TEAL),
    "louis": ("Ruihe “Louis” Zhang", AMBER),
    "othmane": ("Othmane Zizi", BLUE),
    "rui": ("Rui Zhao", PURPLE),
}

EMU_IN = 914400
SW, SH = 13.333, 7.5  # slide size in inches

# --------------------------------------------------------------------------- #
# Asset paths
# --------------------------------------------------------------------------- #
A = {
    "eda": CHARTS / "eda_demand_heatmap.png",
    "results": PRES / "results_performance_slide.png",
    "rebalance": CHARTS / "rebalancing_priority.png",
    "occupancy": CHARTS / "occupancy_trajectory.png",
    "map": CHARTS / "netflow_station_map.png",
    "shap_dep": PRES / "shap_beeswarm_departure.png",
    "shap_arr": PRES / "shap_beeswarm_arrival.png",
    "mlflow_dep": PRES / "mlflow_departure_runs.png",
    "mlflow_arr": PRES / "mlflow_arrival_runs.png",
    "mlflow_reg": PRES / "mlflow_model_registry.png",
    "mlflow_model": PRES / "mlflow_model_departure.png",
    "drift_feat": PRES / "drift_feature_departure_oct.png",
    "drift_concept": PRES / "drift_concept_departure_oct.png",
    "ec2": PRES / "ec2_streamlit_container_running.png",
    "st_7day": CHARTS / "streamlit_7day_forecast.png",
    "st_rebal": CHARTS / "streamlit_rebalancing.png",
    "st_custom": CHARTS / "streamlit_custom_inputs.png",
    "st_monitor": CHARTS / "streamlit_monitoring.png",
}

REPO_URL = "github.com/ruizhaoca/bixi-demand-mlops-platform"
CLOUD_URL = "bixidemandlocal.streamlit.app"
EC2_URL = "3.16.250.166:8501"


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def _set_font(run, size, bold=False, color=INK, name=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name
    run.font.color.rgb = color


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, bold=False, color=INK, name=FONT, align=PP_ALIGN.LEFT,
         space_after=4, space_before=0, italic=False, first=False, level=0,
         bullet=None, line_spacing=None):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if bullet:
        run = p.add_run()
        run.text = f"{bullet}  "
        _set_font(run, size, bold=True, color=RED, name=name)
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold=bold, color=color, name=name, italic=italic)
    return p


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         shadow=False, radius=0.08):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        # lightweight outer shadow
    try:
        if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def node(slide, l, t, w, h, title, sub=None, fill=PALER, line=LGREY, title_color=INK,
         title_size=12, sub_size=9.5, line_w=1.0, radius=0.10, title_bold=True):
    sp = rect(slide, l, t, w, h, fill=fill, line=line, line_w=line_w, radius=radius)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    para(tf, title, title_size, bold=title_bold, color=title_color, align=PP_ALIGN.CENTER,
         space_after=0, first=True)
    if sub:
        para(tf, sub, sub_size, color=DARK, align=PP_ALIGN.CENTER, space_after=0, space_before=1)
    return sp


def connect(slide, x1, y1, x2, y2, color=GREY, w=1.5, arrow=True, kind=MSO_CONNECTOR.STRAIGHT):
    cn = slide.shapes.add_connector(kind, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    if arrow:
        ln = cn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return cn


def img_fit(slide, path, l, t, w, h, align="center", valign="middle", frame=False,
            frame_color=RGBColor(0xE5, 0xE7, 0xEB)):
    """Place an image fit (contain) into the (l,t,w,h) box, preserving aspect."""
    path = Path(path)
    if not path.exists():
        ph = rect(slide, l, t, w, h, fill=PALE, line=LGREY, shape=MSO_SHAPE.RECTANGLE)
        tf = ph.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, f"[missing: {path.name}]", 11, color=GREY, align=PP_ALIGN.CENTER, first=True)
        return ph
    iw, ih = Image.open(path).size
    box_ar = w / h
    img_ar = iw / ih
    if img_ar > box_ar:
        nw = w
        nh = w / img_ar
    else:
        nh = h
        nw = h * img_ar
    nl = l + {"center": (w - nw) / 2, "left": 0, "right": w - nw}[align]
    nt = t + {"middle": (h - nh) / 2, "top": 0, "bottom": h - nh}[valign]
    if frame:
        rect(slide, nl, nt, nw, nh, fill=None, line=frame_color, line_w=1.0,
             shape=MSO_SHAPE.RECTANGLE)
    return slide.shapes.add_picture(str(path), Inches(nl), Inches(nt), Inches(nw), Inches(nh))


# --------------------------------------------------------------------------- #
# Slide chrome
# --------------------------------------------------------------------------- #
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def chrome(slide, presenter_key, title, kicker=None):
    """Accent bar (presenter colour) + title + presenter chip."""
    _, color = PRESENTER[presenter_key]
    # top accent bar
    rect(slide, 0, 0, SW, 0.16, fill=color, shape=MSO_SHAPE.RECTANGLE)
    # title
    _, tf = textbox(slide, 0.55, 0.42, 9.7, 1.05)
    if kicker:
        para(tf, kicker.upper(), 11.5, bold=True, color=color, space_after=2, first=True)
        para(tf, title, 27, bold=True, color=INK, space_after=0)
    else:
        para(tf, title, 28, bold=True, color=INK, space_after=0, first=True)
    # presenter chip (top-right)
    name = PRESENTER[presenter_key][0]
    chip = rect(slide, SW - 3.05, 0.42, 2.5, 0.46, fill=color, radius=0.5)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(ctf, name, 11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True, space_after=0)
    return color


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def url_pill(slide, l, t, label, url, color=RED, w=None):
    w = w or (0.18 * 0.1 * (len(label) + len(url)) + 1.6)
    sp = rect(slide, l, t, w, 0.42, fill=PALER, line=color, line_w=1.0, radius=0.5)
    tf = sp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{label}  "
    _set_font(r, 10.5, bold=True, color=color)
    r = p.add_run(); r.text = url
    _set_font(r, 10.5, bold=False, color=DARK, name=MONO)
    return sp


# --------------------------------------------------------------------------- #
# SLIDES
# --------------------------------------------------------------------------- #
def s_title(prs):
    slide = blank(prs)
    rect(slide, 0, 0, SW, SH, fill=INK, shape=MSO_SHAPE.RECTANGLE)
    rect(slide, 0, 0, SW, 0.28, fill=RED, shape=MSO_SHAPE.RECTANGLE)
    rect(slide, 0, SH - 0.18, SW, 0.18, fill=RED, shape=MSO_SHAPE.RECTANGLE)

    _, tf = textbox(slide, 0.9, 1.15, 11.5, 2.6)
    para(tf, "THE BIXI CREW", 20, bold=True, color=RGBColor(0xFC, 0xA5, 0xA5),
         space_after=8, first=True)
    para(tf, "BIXI Demand MLOps Platform", 40, bold=True, color=WHITE, space_after=6)
    para(tf, "15-minute demand forecasting for every Montreal station — "
            "departures & arrivals — productionised on AWS",
         16, color=RGBColor(0xD1, 0xD5, 0xDB), space_after=0)

    # member table
    members = [
        ("Othmane Zizi", "othmane-zizi-pro"),
        ("Sarah Liu", "sarahliu-mma"),
        ("Ruihe “Louis” Zhang", "Mudkipython"),
        ("Rui Zhao", "ruizhaoca"),
    ]
    x = 0.9
    cw = 2.95
    y = 4.15
    for nm, gh in members:
        card = rect(slide, x, y, cw - 0.18, 1.05, fill=RGBColor(0x1F, 0x29, 0x37),
                    line=RGBColor(0x37, 0x41, 0x51), line_w=1.0, radius=0.12)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        para(tf, nm, 13, bold=True, color=WHITE, space_after=2, first=True)
        p = tf.paragraphs[0]
        p2 = para(tf, f"@{gh}", 11, color=RGBColor(0xFC, 0xA5, 0xA5), name=MONO, space_after=0)
        x += cw

    _, tf = textbox(slide, 0.9, 5.55, 11.5, 1.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Repository   "
    _set_font(r, 12.5, bold=True, color=RGBColor(0x9C, 0xA3, 0xAF))
    r = p.add_run(); r.text = REPO_URL
    _set_font(r, 12.5, bold=False, color=WHITE, name=MONO)
    p = para(tf, "", 6, space_after=0)
    p = tf.add_paragraph()
    r = p.add_run(); r.text = "Live demo   "
    _set_font(r, 12.5, bold=True, color=RGBColor(0x9C, 0xA3, 0xAF))
    r = p.add_run(); r.text = CLOUD_URL
    _set_font(r, 12.5, color=WHITE, name=MONO)
    r = p.add_run(); r.text = "      EC2   "
    _set_font(r, 12.5, bold=True, color=RGBColor(0x9C, 0xA3, 0xAF))
    r = p.add_run(); r.text = EC2_URL
    _set_font(r, 12.5, color=WHITE, name=MONO)

    notes(slide, NOTES[0])


def s_problem(prs):
    slide = blank(prs)
    chrome(slide, "sarah", "The rebalancing problem — and why our design fits it",
           kicker="Introduction & business value")
    # left: problem; right: design choices
    _, tf = textbox(slide, 0.55, 1.75, 6.0, 5.2)
    para(tf, "The operational pain", 15, bold=True, color=TEAL, space_after=6, first=True)
    for b in [
        "BIXI runs ~1,100+ stations. Bikes pile up downtown and run out in "
        "residential areas — riders find empty docks or no free dock.",
        "Operators must move bikes by truck before pressure builds. That needs "
        "demand foresight at operational resolution, not daily totals.",
        "A course-1 prototype predicted hourly demand for the top ~400 stations "
        "— too coarse and too partial to act on.",
    ]:
        para(tf, b, 13, color=DARK, bullet="●", space_after=9, line_spacing=1.05)

    _, tf = textbox(slide, 6.95, 1.75, 5.8, 5.2)
    para(tf, "Our design choices", 15, bold=True, color=TEAL, space_after=6, first=True)
    for head, body in [
        ("15-minute resolution", "4× finer than hourly — matches how fast a station empties."),
        ("Departures vs arrivals, separately", "one shared pipeline, run twice — the gap is the rebalancing signal."),
        ("All ~1,100 stations", "every station scored, not just the busy core."),
        ("Leakage-safe & productionised", "honest evaluation, then served live as an MLOps platform."),
    ]:
        p = para(tf, head, 13, bold=True, color=INK, bullet="▸", space_after=1)
        para(tf, body, 11.5, color=GREY, space_after=8, line_spacing=1.02)
    notes(slide, NOTES[1])


def s_data(prs):
    slide = blank(prs)
    chrome(slide, "sarah", "Data: sources, cleaning & what the demand looks like",
           kicker="Data engineering")
    # flow boxes
    y = 1.7
    node(slide, 0.55, y, 2.55, 0.95, "BIXI open data",
         "2024 + May/Oct 2025 trips, all stations", fill=PALER, line=TEAL, title_size=12)
    node(slide, 3.35, y, 2.5, 0.95, "Open-Meteo weather",
         "15-min temp, precip, wind, humidity", fill=PALER, line=TEAL, title_size=12)
    connect(slide, 3.1, y + 0.47, 3.35, y + 0.47, color=LGREY)
    node(slide, 6.1, y, 2.7, 0.95, "Clean & reshape",
         "trips → 15-min station demand; dep/arr split", fill=PALER, line=TEAL, title_size=12)
    connect(slide, 5.85, y + 0.47, 6.1, y + 0.47, color=LGREY)
    node(slide, 9.05, y, 3.7, 0.95, "Tidy 15-min demand table",
         "~1,100 stations × 96 daily slots × day-of-week", fill=RGBColor(0xEC, 0xFD, 0xF5),
         line=TEAL, title_size=12)
    connect(slide, 8.8, y + 0.47, 9.05, y + 0.47, color=LGREY)
    # heatmap
    img_fit(slide, A["eda"], 0.55, 2.95, 12.25, 3.9, frame=True)
    notes(slide, NOTES[2])


def s_features(prs):
    slide = blank(prs)
    chrome(slide, "sarah", "Feature engineering — informative and leakage-safe",
           kicker="Features & honest evaluation")
    # left features
    _, tf = textbox(slide, 0.55, 1.75, 6.1, 5.2)
    para(tf, "Engineered features", 15, bold=True, color=TEAL, space_after=6, first=True)
    for head, body in [
        ("Cyclical time", "slot-of-day (sin/cos), day-of-week, month — no false ordering."),
        ("2024 profile baselines", "historical average demand + prev-15min / prev-1h / "
            "yesterday-same-slot, built leave-one-out on training rows."),
        ("Weather joins", "15-minute temperature, precipitation, wind, humidity, code."),
        ("Advanced station encoding", "frequency + smoothed target encoding of the "
            "high-cardinality station name."),
        ("Demand tiers", "low / medium / high — used later for fairness parity."),
    ]:
        p = para(tf, head, 12.5, bold=True, color=INK, bullet="▸", space_after=1)
        para(tf, body, 11, color=GREY, space_after=6, line_spacing=1.02)

    # right leakage-safety
    rect(slide, 6.95, 1.75, 5.85, 2.55, fill=PALER, line=TEAL, line_w=1.25, radius=0.06)
    _, tf = textbox(slide, 7.2, 1.95, 5.35, 2.2)
    para(tf, "Leakage-safe by construction", 14, bold=True, color=TEAL, space_after=6, first=True)
    for b in [
        "Strict temporal split — train 2024, validate May-2025, test Oct-2025.",
        "Encoders & baselines fit on TRAIN only, then applied forward.",
        "Leave-one-out baselines so a row never sees its own label.",
    ]:
        para(tf, b, 11.5, color=DARK, bullet="✓", space_after=5, line_spacing=1.02)

    # timeline strip
    y = 4.75
    node(slide, 6.95, y, 1.85, 0.9, "TRAIN", "2024 (full year)", fill=RGBColor(0xEC, 0xFD, 0xF5),
         line=TEAL, title_size=13)
    node(slide, 8.95, y, 1.85, 0.9, "VALIDATION", "May 2025", fill=PALER, line=LGREY, title_size=13)
    node(slide, 10.95, y, 1.85, 0.9, "TEST", "Oct 2025", fill=PALER, line=LGREY, title_size=13)
    connect(slide, 8.8, y + 0.45, 8.95, y + 0.45, color=LGREY)
    connect(slide, 10.8, y + 0.45, 10.95, y + 0.45, color=LGREY)
    _, tf = textbox(slide, 6.95, y + 1.0, 5.85, 0.5)
    para(tf, "Time →   encoders & baselines learned on the left, only applied to the right",
         10.5, color=GREY, italic=True, first=True)
    notes(slide, NOTES[3])


def s_architecture(prs):
    slide = blank(prs)
    chrome(slide, "louis", "Architecture — GitHub to AWS, all infrastructure-as-code",
           kicker="AWS & productionised structure")
    # CI band
    node(slide, 0.55, 1.7, 2.2, 0.85, "GitHub", "main + PRs", fill=PALER, line=AMBER, title_size=12)
    node(slide, 3.15, 1.7, 2.7, 0.85, "GitHub Actions CI", "pytest + build images",
         fill=PALER, line=AMBER, title_size=12)
    node(slide, 6.25, 1.7, 2.6, 0.85, "Docker images", "train + Streamlit serving",
         fill=PALER, line=AMBER, title_size=12)
    connect(slide, 2.75, 2.12, 3.15, 2.12, color=AMBER)
    connect(slide, 5.85, 2.12, 6.25, 2.12, color=AMBER)

    # AWS container
    aws = rect(slide, 0.55, 2.95, 8.3, 2.55, fill=RGBColor(0xFF, 0xFB, 0xEB), line=AMBER,
               line_w=1.5, radius=0.04)
    _, tf = textbox(slide, 0.72, 3.05, 7.9, 0.4)
    para(tf, "AWS  us-east-2  ·  provisioned by AWS CDK (infra/)", 12.5, bold=True,
         color=AMBER, first=True)
    stacks = [
        ("BixiNetwork", "VPC, public subnets"),
        ("BixiStorage", "S3 bucket + SSM"),
        ("BixiMlflow", "MLflow on EC2 + S3"),
        ("BixiBatch", "ECR + AWS Batch"),
    ]
    sx = 0.72
    for nm, sub in stacks:
        node(slide, sx, 3.55, 1.92, 0.85, nm, sub, fill=WHITE, line=AMBER, title_size=11.5,
             sub_size=9)
        sx += 2.0
    _, tf = textbox(slide, 0.72, 4.55, 7.9, 0.85)
    p = para(tf, "AWS Batch runs  ", 11.5, color=DARK, first=True, space_after=2)
    r = p.add_run(); r.text = "python -m bixi.pipeline"
    _set_font(r, 11.5, bold=True, color=INK, name=MONO)
    r = p.add_run(); r.text = "  over the full dataset"
    _set_font(r, 11.5, color=DARK)
    para(tf, "source data  s3://insy684   →   checkpoints, models, reports → S3 + MLflow",
         11, color=GREY, name=MONO, space_after=0)
    connect(slide, 7.55, 2.55, 7.55, 3.55, color=AMBER)  # docker -> batch

    # serving
    serv = rect(slide, 9.15, 1.7, 3.65, 3.8, fill=PALER, line=PURPLE, line_w=1.25, radius=0.05)
    _, tf = textbox(slide, 9.32, 1.8, 3.3, 0.4)
    para(tf, "Serving", 12.5, bold=True, color=PURPLE, first=True)
    node(slide, 9.32, 2.3, 3.3, 1.0, "Streamlit Community Cloud",
         "app.py · committed artifacts, no AWS at runtime", fill=WHITE, line=PURPLE,
         title_size=11.5, sub_size=9)
    node(slide, 9.32, 3.45, 3.3, 1.0, "EC2 Streamlit container",
         "app_ec2.py · same artifacts from S3", fill=WHITE, line=PURPLE,
         title_size=11.5, sub_size=9)
    _, tf = textbox(slide, 9.32, 4.6, 3.3, 0.8)
    para(tf, CLOUD_URL, 9.5, color=DARK, name=MONO, first=True, space_after=2)
    para(tf, EC2_URL, 9.5, color=DARK, name=MONO, space_after=0)
    connect(slide, 8.85, 3.6, 9.32, 2.8, color=LGREY)  # models -> serving

    _, tf = textbox(slide, 0.55, 5.75, 12.3, 0.9)
    para(slide_tf := tf, "Identical code runs locally on a station subsample and on AWS Batch over "
         "the full dataset. No credentials in git — default boto3 chain (SSO locally, IAM role in cloud).",
         12, color=DARK, italic=True, first=True)
    notes(slide, NOTES[4])


def s_pipeline(prs):
    slide = blank(prs)
    chrome(slide, "louis", "A resumable, staged pipeline — same code everywhere",
           kicker="Productionisation")
    stages = [
        ("ingest", "trips + weather → 15-min demand"),
        ("features", "leakage-safe feature tables"),
        ("data", "encode, tiers, range-filter"),
        ("train", "candidates + FLAML + Optuna"),
        ("explain", "SHAP + LIME"),
        ("fairness", "error parity"),
        ("drift", "Evidently 4-type"),
        ("register", "promote to production"),
    ]
    # two rows of 4
    bw, bh = 2.85, 1.15
    gx, gy = 0.18, 0.7
    x0, y0 = 0.62, 1.85
    positions = []
    for i, (nm, sub) in enumerate(stages):
        row, col = divmod(i, 4)
        x = x0 + col * (bw + gx)
        y = y0 + row * (bh + gy)
        positions.append((x, y))
        fill = RGBColor(0xFF, 0xFB, 0xEB) if i >= 5 else WHITE
        node(slide, x, y, bw, bh, nm, sub, fill=fill, line=AMBER, title_size=14, sub_size=10,
             radius=0.12)
    # arrows within row 1
    for col in range(3):
        x = x0 + col * (bw + gx)
        connect(slide, x + bw, y0 + bh / 2, x + bw + gx, y0 + bh / 2, color=AMBER)
    # arrow wrap row1->row2 (down on right, along)
    connect(slide, positions[3][0] + bw / 2, y0 + bh, positions[4][0] + bw / 2,
            y0 + bh + gy, color=AMBER)
    for col in range(3):
        x = x0 + col * (bw + gx)
        y = y0 + bh + gy
        connect(slide, x + bw, y + bh / 2, x + bw + gx, y + bh / 2, color=AMBER)

    # success marker note
    rect(slide, 0.62, 5.55, 12.1, 1.15, fill=PALER, line=AMBER, line_w=1.0, radius=0.05)
    _, tf = textbox(slide, 0.85, 5.7, 11.7, 0.95)
    p = para(tf, "Each stage writes a  ", 12, color=DARK, first=True, space_after=3)
    r = p.add_run(); r.text = "_SUCCESS"
    _set_font(r, 12, bold=True, color=INK, name=MONO)
    r = p.add_run(); r.text = "  marker to S3 — a run resumes from any step. Default run starts at "
    _set_font(r, 12, color=DARK)
    r = p.add_run(); r.text = "data"
    _set_font(r, 12, bold=True, color=INK, name=MONO)
    r = p.add_run(); r.text = " (cleaned data already in S3); full rebuild starts at "
    _set_font(r, 12, color=DARK)
    r = p.add_run(); r.text = "ingest"
    _set_font(r, 12, bold=True, color=INK, name=MONO)
    r = p.add_run(); r.text = "."
    _set_font(r, 12, color=DARK)
    para(tf, "CI builds both Docker images and smoke-tests the pipeline image with no AWS.",
         11.5, color=GREY, italic=True, space_after=0)
    notes(slide, NOTES[5])


def s_repo(prs):
    slide = blank(prs)
    chrome(slide, "louis", "Repository structure — single-responsibility modules",
           kicker="Clean, config-driven design")
    tree = (
        "src/bixi/                  the pipeline package\n"
        " ├ config.py             central config + data/feature contract + stages\n"
        " ├ io.py                 S3 + local I/O (default boto3 chain)\n"
        " ├ ingest.py / *_cleaning  raw trips + weather → 15-min demand\n"
        " ├ feature_engineering   leakage-safe feature tables\n"
        " ├ data.py               range filter, station encoding, tiers\n"
        " ├ models.py             candidates, FLAML, Optuna, metrics\n"
        " ├ pipeline.py           resumable staged runner\n"
        " ├ explain / fairness / drift   SHAP+LIME, parity, Evidently\n"
        " ├ registry.py           MLflow tracking + promotion\n"
        " ├ rebalancing.py        net-flow priorities (dep+arr)\n"
        " └ streamlit_*_serving   local + S3 serving helpers\n"
        "app.py / app_ec2.py        Streamlit UIs (Cloud / EC2)\n"
        "infra/                     AWS CDK app (4 stacks)\n"
        "docker/                    Dockerfile.train, Dockerfile.streamlit_ec2\n"
        "tests/                     pytest (synthetic data, no network)\n"
        "docs/  scripts/  notebooks/"
    )
    box = rect(slide, 0.55, 1.7, 7.95, 5.25, fill=RGBColor(0x0F, 0x17, 0x2A), line=None, radius=0.03)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_top = Inches(0.16)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    import re
    first = True
    for line in tree.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = Pt(2.5)
        first = False
        # split the path token from its description on the first run of 2+ spaces
        m = re.split(r"\s{2,}", line, maxsplit=1)
        path_txt = m[0]
        comment = m[1].strip() if len(m) > 1 else ""
        r = p.add_run(); r.text = path_txt + ("   " if comment else "")
        _set_font(r, 11, color=RGBColor(0xF8, 0xFA, 0xFC), name=MONO)
        if comment:
            r = p.add_run(); r.text = comment
            _set_font(r, 9.5, color=RGBColor(0x86, 0xEF, 0xAC), name=MONO, italic=True)

    # callouts
    _, tf = textbox(slide, 8.8, 1.75, 4.0, 5.1)
    para(tf, "Why it scores", 15, bold=True, color=AMBER, space_after=8, first=True)
    for head, body in [
        ("Single responsibility", "each module owns one stage of the contract in config.py."),
        ("Config-driven", "env vars switch local subsample ↔ full cloud run — same code."),
        ("Infra as code", "every AWS resource lives in infra/ CDK — nothing clicked."),
        ("Tested & CI-gated", "pytest on synthetic data; PRs must stay green."),
        ("Reproducible", "Docker pins the runtime for training and serving."),
    ]:
        p = para(tf, head, 12.5, bold=True, color=INK, bullet="▸", space_after=1)
        para(tf, body, 11, color=GREY, space_after=7, line_spacing=1.03)
    notes(slide, NOTES[6])


def s_modeling(prs):
    slide = blank(prs)
    chrome(slide, "othmane", "Modelling — a multi-model search, auto-selected",
           kicker="Predictive modelling")
    # funnel of boxes
    y = 1.95
    node(slide, 0.55, y, 2.75, 1.5,
         "Candidates", "LightGBM  +  XGBoost\nbaselines", fill=PALER, line=BLUE,
         title_size=14, sub_size=11)
    node(slide, 0.55, y + 1.75, 2.75, 1.2, "FLAML AutoML",
         "automated model + config search", fill=PALER, line=BLUE, title_size=13, sub_size=10.5)
    node(slide, 3.6, y + 0.85, 2.75, 1.4, "Optuna HPO",
         "Bayesian search\nevery trial → MLflow", fill=RGBColor(0xEF, 0xF6, 0xFF), line=BLUE,
         title_size=14, sub_size=11)
    connect(slide, 3.3, y + 0.75, 3.6, y + 1.2, color=BLUE)
    connect(slide, 3.3, y + 2.35, 3.6, y + 1.9, color=BLUE)
    node(slide, 6.65, y + 0.95, 2.85, 1.2, "Select best", "by validation RMSE",
         fill=RGBColor(0xEF, 0xF6, 0xFF), line=BLUE, title_size=14, sub_size=11)
    connect(slide, 6.35, y + 1.55, 6.65, y + 1.55, color=BLUE)
    node(slide, 9.8, y + 0.95, 2.95, 1.2, "lgbm_optuna", "promoted → production alias",
         fill=RGBColor(0xDC, 0xFC, 0xE7), line=RGBColor(0x16, 0xA3, 0x4A), title_size=15,
         sub_size=10.5, title_color=RGBColor(0x15, 0x80, 0x3D))
    connect(slide, 9.5, y + 1.55, 9.8, y + 1.55, color=BLUE)

    rect(slide, 0.55, 5.75, 12.25, 1.1, fill=PALER, line=BLUE, line_w=1.0, radius=0.05)
    _, tf = textbox(slide, 0.8, 5.9, 11.8, 0.9)
    para(tf, "Run once per target (departure & arrival). Both targets independently select "
            "lgbm_optuna — a consistent, automatic, reproducible winner, not a hand-picked model.",
         12.5, color=DARK, first=True, line_spacing=1.05)
    notes(slide, NOTES[7])


def s_mlflow(prs):
    slide = blank(prs)
    chrome(slide, "othmane", "MLflow — every run tracked, best model registered",
           kicker="Experiment tracking & registry")
    img_fit(slide, A["mlflow_dep"], 0.55, 1.75, 4.05, 3.05, frame=True)
    img_fit(slide, A["mlflow_arr"], 4.72, 1.75, 4.05, 3.05, frame=True)
    img_fit(slide, A["mlflow_reg"], 8.9, 1.75, 3.9, 3.05, frame=True)
    _, tf = textbox(slide, 0.55, 1.5, 12, 0.3)
    # captions
    for x, txt in [(0.55, "departure · 73 runs"), (4.72, "arrival · 57 runs"),
                   (8.9, "model registry · production alias")]:
        _, tfc = textbox(slide, x, 4.85, 4.0, 0.3)
        para(tfc, txt, 10.5, bold=True, color=GREY, align=PP_ALIGN.CENTER, first=True)

    _, tf = textbox(slide, 0.55, 5.35, 12.25, 1.4)
    para(tf, "Tracking server on EC2 with an S3 artifact store. Experiments "
            "bixi-demand-departure / -arrival capture every candidate, FLAML and Optuna run "
            "(params, metrics, model). The best run per target is registered and promoted to the "
            "production alias — the exact artifact the apps serve.", 12.5, color=DARK,
         first=True, line_spacing=1.08)
    notes(slide, NOTES[8])


def s_results(prs):
    slide = blank(prs)
    # results image is a full 16:9 slide already; place near full-bleed with a chip
    img_fit(slide, A["results"], 0.0, 0.0, SW, SH)
    rect(slide, 0, 0, SW, 0.16, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
    chip = rect(slide, SW - 3.05, 0.30, 2.5, 0.42, fill=BLUE, radius=0.5)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(ctf, PRESENTER["othmane"][0], 11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         first=True, space_after=0)
    notes(slide, NOTES[9])


def s_explain(prs):
    slide = blank(prs)
    chrome(slide, "othmane", "Explainability — SHAP & LIME on both models",
           kicker="Responsible AI")
    img_fit(slide, A["shap_dep"], 0.55, 1.7, 5.95, 3.5, frame=True)
    img_fit(slide, A["shap_arr"], 6.65, 1.7, 5.95, 3.5, frame=True)
    for x, txt in [(0.55, "SHAP · departure"), (6.65, "SHAP · arrival")]:
        _, tfc = textbox(slide, x, 5.25, 5.95, 0.3)
        para(tfc, txt, 10.5, bold=True, color=GREY, align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(slide, 0.55, 5.65, 12.25, 1.2)
    para(tf, "Most signal comes from the 2024 historical baselines and cyclical time-of-day; "
            "weather is a secondary driver. LIME force plots explain individual predictions. "
            "Same story for departures and arrivals — the model is behaving sensibly.",
         12.5, color=DARK, first=True, line_spacing=1.08)
    notes(slide, NOTES[10])


def s_fairness_drift(prs):
    slide = blank(prs)
    chrome(slide, "othmane", "Fairness & four-type drift monitoring",
           kicker="Responsible AI & monitoring")
    # left fairness text
    rect(slide, 0.55, 1.7, 4.0, 5.0, fill=PALER, line=BLUE, line_w=1.0, radius=0.04)
    _, tf = textbox(slide, 0.78, 1.85, 3.55, 4.7)
    para(tf, "Fairness", 15, bold=True, color=BLUE, space_after=6, first=True)
    for b in [
        "Error parity (RMSE/MAE) across demand tiers and geography.",
        "Accuracy concentrates in high-demand stations; low-tier R² is near zero.",
        "Flagged so operators don’t over-trust quiet-station forecasts.",
    ]:
        para(tf, b, 11.5, color=DARK, bullet="●", space_after=8, line_spacing=1.05)
    para(tf, "Four drift types (Evidently)", 14, bold=True, color=BLUE, space_after=5,
         space_before=6)
    for b in ["feature", "target", "prediction", "concept"]:
        para(tf, b, 11.5, color=DARK, bullet="▸", space_after=4)

    img_fit(slide, A["drift_feat"], 4.75, 1.7, 4.0, 4.5, frame=True)
    img_fit(slide, A["drift_concept"], 8.85, 1.7, 3.95, 4.5, frame=True)
    for x, txt, w in [(4.75, "feature drift · Oct-2025", 4.0),
                      (8.85, "concept drift · Oct-2025", 3.95)]:
        _, tfc = textbox(slide, x, 6.25, w, 0.3)
        para(tfc, txt, 10.5, bold=True, color=GREY, align=PP_ALIGN.CENTER, first=True)
    notes(slide, NOTES[11])


def s_rebalancing(prs):
    slide = blank(prs)
    chrome(slide, "rui", "Net-flow rebalancing — turning forecasts into action",
           kicker="From prediction to operations")
    _, tf = textbox(slide, 0.55, 1.62, 12.25, 0.55)
    p = para(tf, "net_flow = arrival − departure", 13, bold=True, color=PURPLE, first=True,
             space_after=0)
    r = p.add_run()
    r.text = "   → cumulate across the day → read peak deficit (needs bikes) / peak surplus (needs docks) → rank by severity"
    _set_font(r, 12, color=DARK)
    img_fit(slide, A["rebalance"], 0.55, 2.25, 6.5, 4.55, frame=True)
    img_fit(slide, A["map"], 7.2, 2.25, 5.6, 4.55, frame=True)
    notes(slide, NOTES[12])


def s_demo(prs):
    slide = blank(prs)
    chrome(slide, "rui", "The Streamlit app — live demo", kicker="Serving & product")
    # 2x2 grid of streamlit screenshots
    cells = [
        (A["st_7day"], "7-Day Forecast"),
        (A["st_rebal"], "Rebalancing Priorities"),
        (A["st_custom"], "Custom Inputs"),
        (A["st_monitor"], "Model Monitoring"),
    ]
    gw, gh = 4.55, 2.35
    x0, y0 = 0.55, 1.72
    gx, gy = 0.2, 0.32
    for i, (path, cap) in enumerate(cells):
        row, col = divmod(i, 2)
        x = x0 + col * (gw + gx)
        y = y0 + row * (gh + gy)
        img_fit(slide, path, x, y, gw, gh, frame=True)
        _, tfc = textbox(slide, x, y + gh + 0.01, gw, 0.26)
        para(tfc, cap, 10, bold=True, color=GREY, align=PP_ALIGN.CENTER, first=True)

    # right rail: ec2 + urls
    img_fit(slide, A["ec2"], 9.95, 1.72, 2.9, 2.0, frame=True)
    _, tfc = textbox(slide, 9.95, 3.74, 2.9, 0.26)
    para(tfc, "EC2 container running", 10, bold=True, color=GREY, align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(slide, 9.95, 4.25, 2.95, 2.4)
    para(tf, "Two live deployments", 13, bold=True, color=PURPLE, space_after=6, first=True)
    para(tf, "Community Cloud", 11, bold=True, color=INK, space_after=1)
    para(tf, CLOUD_URL, 10, color=DARK, name=MONO, space_after=8)
    para(tf, "EC2 (S3-backed)", 11, bold=True, color=INK, space_after=1)
    para(tf, EC2_URL, 10, color=DARK, name=MONO, space_after=8)
    para(tf, "Four pages · one shared UI · no AWS needed for the Cloud app.",
         10.5, color=GREY, italic=True, space_after=0, line_spacing=1.05)
    notes(slide, NOTES[13])


def s_conclusion(prs):
    slide = blank(prs)
    chrome(slide, "rui", "What it means — value, limits, and what’s next",
           kicker="Conclusion")
    img_fit(slide, A["occupancy"], 0.55, 1.7, 6.4, 4.0, frame=True)
    _, tfc = textbox(slide, 0.55, 5.72, 6.4, 0.3)
    para(tfc, "relative occupancy starts from a common zero reference", 10, color=GREY,
         italic=True, align=PP_ALIGN.CENTER, first=True)

    _, tf = textbox(slide, 7.2, 1.7, 5.6, 5.1)
    para(tf, "Operational value", 14, bold=True, color=PURPLE, space_after=4, first=True)
    para(tf, "A daily, ranked priority list of which stations to service first — grounded "
            "in 15-minute departure & arrival forecasts.", 11.5, color=DARK, space_after=10,
         line_spacing=1.05)
    para(tf, "Honest limitation", 14, bold=True, color=PURPLE, space_after=4)
    para(tf, "No dock capacity or real-time occupancy in the trip data, so the trajectory is a "
            "relative risk ranking — not exact stockout clock-times or fill levels.",
         11.5, color=DARK, space_after=10, line_spacing=1.05)
    para(tf, "MLOps maturity & roadmap", 14, bold=True, color=PURPLE, space_after=4)
    for b in [
        "Tracking + registry, explainability, fairness, 4-type drift, CI, Docker, CDK — end to end.",
        "Next: live station-status feed for absolute occupancy; scheduled drift-triggered retrains.",
    ]:
        para(tf, b, 11.5, color=DARK, bullet="▸", space_after=6, line_spacing=1.04)
    notes(slide, NOTES[14])


def s_thanks(prs):
    slide = blank(prs)
    rect(slide, 0, 0, SW, SH, fill=INK, shape=MSO_SHAPE.RECTANGLE)
    rect(slide, 0, 0, SW, 0.28, fill=RED, shape=MSO_SHAPE.RECTANGLE)
    rect(slide, 0, SH - 0.18, SW, 0.18, fill=RED, shape=MSO_SHAPE.RECTANGLE)
    _, tf = textbox(slide, 0.9, 1.5, 11.5, 2.0)
    para(tf, "Thank you", 46, bold=True, color=WHITE, first=True, space_after=6)
    para(tf, "Questions & live demo welcome", 18, color=RGBColor(0xD1, 0xD5, 0xDB), space_after=0)

    _, tf = textbox(slide, 0.9, 3.95, 11.5, 2.6)
    rows = [
        ("Repository", REPO_URL),
        ("Live demo (Community Cloud)", CLOUD_URL),
        ("Live demo (EC2)", EC2_URL),
    ]
    for label, url in rows:
        p = tf.paragraphs[0] if (label == "Repository") else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run(); r.text = f"{label}    "
        _set_font(r, 14, bold=True, color=RGBColor(0x9C, 0xA3, 0xAF))
        r = p.add_run(); r.text = url
        _set_font(r, 14, color=WHITE, name=MONO)
    p = tf.add_paragraph(); p.space_before = Pt(6)
    r = p.add_run()
    r.text = "The Bixi Crew  ·  Othmane Zizi  ·  Sarah Liu  ·  Ruihe “Louis” Zhang  ·  Rui Zhao"
    _set_font(r, 13, bold=True, color=RGBColor(0xFC, 0xA5, 0xA5))
    notes(slide, NOTES[15])


# --------------------------------------------------------------------------- #
# Speaker notes (also written to docs/presentation/speaker_notes.tex source)
# --------------------------------------------------------------------------- #
NOTES = [
    # 0 title
    "Shared open (Sarah leads). We are The Bixi Crew. Our project forecasts 15-minute "
    "BIXI demand for every Montreal station, separately for departures and arrivals, and turns "
    "that into a rebalancing tool — productionised end-to-end on AWS. Four of us each own a "
    "part: I introduce the problem and the data, Louis covers the AWS structure, Othmane the "
    "modelling, Rui the rebalancing layer and the live app. Repo and two live demos are on screen.",
    # 1 problem
    "Sarah (1/3). BIXI has ~1,100+ stations. Bikes pile up downtown and vanish in residential "
    "areas, so riders hit empty or full docks. Operators rebalance with trucks but need foresight "
    "at operational resolution. The course-1 prototype was hourly and only the top ~400 stations "
    "— too coarse and partial. So we made four deliberate design choices: 15-minute resolution "
    "(4× finer), departures and arrivals predicted separately (their gap is the signal), all "
    "stations, and a leakage-safe, productionised platform.",
    # 2 data
    "Sarah (2/3). Sources: BIXI open-data trips for 2024 plus May and October 2025, joined to "
    "Open-Meteo 15-minute weather. We clean trips into a tidy 15-minute station-demand table, split "
    "into departures and arrivals. The heatmap is average demand by weekday and 15-minute slot over "
    "all stations: quiet overnight, building through midday, and busiest on weekday evenings. That "
    "structure is exactly what our time features capture.",
    # 3 features
    "Sarah (3/3). Features: cyclical slot-of-day, day-of-week and month; 2024 historical baselines "
    "(average plus recent-lag baselines) built leave-one-out; 15-minute weather; and advanced "
    "frequency + smoothed target encoding of the station name. Crucially everything is leakage-safe: "
    "a strict temporal split — train on 2024, validate on May-2025, test on Oct-2025 — with "
    "encoders and baselines fit on train only and applied forward. Over to Louis.",
    # 4 architecture
    "Louis (1/3). Everything is infrastructure-as-code. Push to GitHub triggers Actions CI: it runs "
    "pytest and builds both Docker images. AWS in us-east-2 is provisioned entirely by AWS CDK in "
    "four stacks — BixiNetwork (VPC), BixiStorage (S3 + SSM), BixiMlflow (MLflow on EC2 + S3), "
    "and BixiBatch (ECR + AWS Batch). AWS Batch runs python -m bixi.pipeline over the full dataset "
    "from s3://insy684. Serving is two Streamlit deployments — Community Cloud and an EC2 "
    "container. No credentials in git: default boto3 chain, SSO locally, IAM role in the cloud.",
    # 5 pipeline
    "Louis (2/3). The pipeline is eight ordered stages: ingest, features, data, train, explain, "
    "fairness, drift, register. It is resumable — each stage writes a _SUCCESS marker to S3, so "
    "a run can restart from any step. The default run starts at data because cleaned data already "
    "lives in S3; a full rebuild starts at ingest. The same code runs on a local subsample and on "
    "AWS Batch, and CI smoke-tests the pipeline image with no AWS.",
    # 6 repo
    "Louis (3/3). The repo mirrors the pipeline: single-responsibility modules in src/bixi, each "
    "owning one stage of the contract defined in config.py. It is config-driven — environment "
    "variables flip between a local subsample and the full cloud run with identical code. infra/ is "
    "the CDK app, docker/ pins the runtime, tests/ runs on synthetic data with no network. Clean, "
    "reproducible, and CI-gated. Othmane takes the modelling.",
    # 7 modeling
    "Othmane (1/5). Rather than hand-pick a model, we run a search: LightGBM and XGBoost baselines, "
    "FLAML AutoML, and Optuna Bayesian hyperparameter tuning, with every trial logged to MLflow. The "
    "best model by validation RMSE is selected and promoted automatically. Run once per target; both "
    "departures and arrivals independently land on lgbm_optuna — a consistent, reproducible "
    "winner.",
    # 8 mlflow
    "Othmane (2/5). MLflow gives full traceability: a tracking server on EC2 with an S3 artifact "
    "store. The departure experiment has 73 runs, arrival 57 — every candidate, FLAML and Optuna "
    "trial with its params, metrics and model. The best run per target is registered and promoted to "
    "the production alias, which is exactly the artifact the apps serve.",
    # 9 results
    "Othmane (3/5). Results per split. 15-minute slot demand is genuinely noisy — many zero "
    "slots — so absolute R² is modest by design, but errors are well under one trip per "
    "slot. The tuned model beats a naive baseline on RMSE and R², it is stable across two unseen "
    "months (no overfitting), and its predictive power is concentrated exactly where it matters — "
    "busy stations.",
    # 10 explain
    "Othmane (4/5). Explainability with SHAP and LIME on both models. The SHAP beeswarms show most "
    "signal comes from the 2024 historical baselines and cyclical time-of-day, with weather as a "
    "secondary driver — the same, sensible story for departures and arrivals. LIME explains "
    "individual predictions on demand.",
    # 11 fairness drift
    "Othmane (5/5). Responsible AI and monitoring. Fairness: we check error parity across demand "
    "tiers and geography — accuracy concentrates in high-demand stations and low-tier R² is "
    "near zero, which we surface so operators don’t over-trust quiet-station forecasts. And we "
    "run all four Evidently drift types — feature, target, prediction and concept — shown "
    "here on October-2025 data. Over to Rui.",
    # 12 rebalancing
    "Rui (1/3). The forecasts are only useful together. Net flow is arrival minus departure per "
    "station per 15-minute slot; cumulating it across the day gives a relative occupancy trajectory. "
    "From it we read each station’s peak deficit — needs bikes, stockout risk — and "
    "peak surplus — needs docks, overflow risk — then rank by severity. The map shows the "
    "pattern: the downtown core fills up and needs docks, the residential periphery drains and needs "
    "bikes.",
    # 13 demo
    "Rui (2/3). This is the live Streamlit app — one shared UI, four pages: a multi-day "
    "15-minute forecast, the rebalancing-priorities map with a ranked list and per-station "
    "trajectory, a custom what-if input page, and a model-monitoring page with the SHAP, fairness and "
    "drift artifacts. It runs on Community Cloud with committed artifacts (no AWS at runtime) and on "
    "an EC2 container backed by S3. Let’s look at it live.",
    # 14 conclusion
    "Rui (3/3). The operational value: a daily ranked list of where to send a rebalancing truck "
    "first. The honest limitation — the trip data has no dock capacity or real-time occupancy, so "
    "the trajectory starts from a common zero and gives a relative ranking, not exact stockout times. "
    "On MLOps maturity we cover the full lifecycle — tracking, registry, explainability, "
    "fairness, four-type drift, CI, Docker and CDK. Next steps: a live station-status feed for "
    "absolute occupancy, and drift-triggered scheduled retrains.",
    # 15 thanks
    "Shared close. Thank you — happy to take questions and walk through the live app. Repo and "
    "both live demos are on screen; the team is The Bixi Crew: Othmane, Sarah, Louis and Rui.",
]


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    s_title(prs)
    s_problem(prs)
    s_data(prs)
    s_features(prs)
    s_architecture(prs)
    s_pipeline(prs)
    s_repo(prs)
    s_modeling(prs)
    s_mlflow(prs)
    s_results(prs)
    s_explain(prs)
    s_fairness_drift(prs)
    s_rebalancing(prs)
    s_demo(prs)
    s_conclusion(prs)
    s_thanks(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(REPO)}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
